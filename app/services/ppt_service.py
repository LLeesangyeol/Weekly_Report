from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import Settings


class PptExtractionError(RuntimeError):
    pass


def _cell_text(cell: Any) -> str:
    return cell.text.strip()


def _normalized_label(value: str) -> str:
    return re.sub(r"\s+", "", value)


def _table_lines(shape: Any) -> list[str]:
    rows = [[_cell_text(cell) for cell in row.cells] for row in shape.table.rows]
    if not rows:
        return []
    label = _normalized_label(rows[0][0]) if rows[0] else ""

    # Common one-page weekly report form: the two right-hand cells are separate
    # work columns, not a single sentence.  Make that structure explicit for the LLM.
    if "금주예정" in label and "업무" in label:
        values = [value for value in rows[0][1:] if value]
        lines = ["[금주 예정 업무]"]
        lines.extend(f"[예정 업무 영역 {index}]\n{value}" for index, value in enumerate(values, 1))
        return lines
    if "주간업무" in label and "실적" in label:
        values = [value for value in rows[0][1:] if value]
        lines = ["[금주 업무 실적]"]
        lines.extend(f"[실적 업무 영역 {index}]\n{value}" for index, value in enumerate(values, 1))
        return lines
    if "주간일정" in label and len(rows) >= 2:
        headers, plans = rows[0][1:], rows[1][1:]
        lines = ["[주간 일정 계획표]"]
        for day, plan in zip(headers, plans):
            if day and plan:
                lines.append(f"- {day}: {plan}")
        return lines
    if label == "부서" and len(rows) >= 2:
        result: list[str] = []
        for row in rows:
            for index in range(0, len(row) - 1, 2):
                key, value = row[index].strip(), row[index + 1].strip()
                if key and value:
                    result.append(f"{key}: {value}")
        return result

    result: list[str] = []
    for row in rows:
        if any(row):
            result.append(" | ".join(row))
    return result


def _shape_text(shape: Any) -> list[str]:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        lines: list[str] = []
        for child in sorted(shape.shapes, key=lambda item: (item.top, item.left)):
            lines.extend(_shape_text(child))
        return lines
    if getattr(shape, "has_table", False):
        return _table_lines(shape)
    if getattr(shape, "has_text_frame", False):
        lines = []
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        return lines
    return []


def extract_pptx_text(path: Path) -> str:
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise PptExtractionError("PPTX 문서를 열 수 없습니다.") from exc
    pages: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        # Spatial ordering keeps top-to-bottom and, for equal rows, left-to-right layout.
        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            lines.extend(_shape_text(shape))
        pages.append(f"--- Slide {number} ---\n" + "\n".join(lines))
    return "\n\n".join(pages).strip()


class PptService:
    def __init__(self, settings: Settings):
        self.settings = settings

    def extract(self, path: Path) -> str:
        if path.suffix.lower() == ".pptx":
            return extract_pptx_text(path)
        if path.suffix.lower() != ".ppt":
            raise PptExtractionError("PPT 또는 PPTX 파일이 아닙니다.")
        return self._convert_and_extract(path)

    def _convert_and_extract(self, source: Path) -> str:
        self.settings.temp_dir.mkdir(parents=True, exist_ok=True)
        # Use a neutral profile path; some Windows LibreOffice builds are sensitive
        # to particular tokens in an isolated profile directory.
        work_dir = Path(tempfile.mkdtemp(prefix="probe-", dir=self.settings.temp_dir))
        try:
            profile_dir = work_dir / "lo-profile"
            profile_dir.mkdir()
            command = [
                self.settings.soffice_path,
                f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(work_dir),
                str(source),
            ]
            try:
                result = subprocess.run(
                    command,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    timeout=self.settings.conversion_timeout_seconds,
                    check=False,
                    shell=False,
                )
                # LibreOffice can occasionally exit during first-time profile setup.
                # Retry exactly once with the same isolated profile and timeout.
                if result.returncode != 0:
                    result = subprocess.run(
                        command,
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                        timeout=self.settings.conversion_timeout_seconds,
                        check=False,
                        shell=False,
                    )
            except FileNotFoundError as exc:
                raise PptExtractionError("LibreOffice(soffice)를 찾을 수 없습니다.") from exc
            except subprocess.TimeoutExpired as exc:
                raise PptExtractionError("PPT 변환 시간이 제한을 초과했습니다.") from exc
            converted = work_dir / f"{source.stem}.pptx"
            if result.returncode != 0 or not converted.is_file():
                detail = (result.stderr or result.stdout or "unknown error").strip()[-500:]
                raise PptExtractionError(f"PPT 변환에 실패했습니다: {detail}")
            return extract_pptx_text(converted)
        finally:
            shutil.rmtree(work_dir, ignore_errors=True)
