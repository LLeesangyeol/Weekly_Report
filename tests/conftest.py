from __future__ import annotations

import io
from collections.abc import Generator
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlalchemy.orm import Session, sessionmaker

from app.config import Settings, get_settings
from app.database import Base, get_db, make_engine
from app.main import app
from app.routers.api import get_report_processor


@pytest.fixture
def settings(tmp_path: Path) -> Settings:
    value = Settings(
        database_url=f"sqlite:///{tmp_path / 'test.db'}",
        upload_dir=tmp_path / "uploads",
        temp_dir=tmp_path / "temp",
        ollama_url="http://ollama.test",
        ollama_model="test-model",
        ollama_timeout_seconds=2,
        max_upload_mb=1,
        max_pdf_pages=10,
        text_chunk_size=5000,
        disk_usage_limit_percent=100,
        min_free_disk_gb=0,
        soffice_path="soffice",
        conversion_timeout_seconds=2,
    )
    value.ensure_directories()
    return value


@pytest.fixture
def session_factory(settings: Settings):
    engine = make_engine(settings.database_url)
    Base.metadata.create_all(engine)
    factory = sessionmaker(bind=engine, expire_on_commit=False, class_=Session)
    try:
        yield factory
    finally:
        engine.dispose()


@pytest.fixture
def db(session_factory) -> Generator[Session, None, None]:
    with session_factory() as session:
        yield session


class NoopProcessor:
    async def process(self, _report_id: int) -> None:
        return None


@pytest.fixture
def client(settings: Settings, session_factory):
    def override_db():
        with session_factory() as session:
            yield session

    app.dependency_overrides[get_db] = override_db
    app.dependency_overrides[get_settings] = lambda: settings
    app.dependency_overrides[get_report_processor] = lambda: NoopProcessor()
    with TestClient(app) as test_client:
        yield test_client
    app.dependency_overrides.clear()


@pytest.fixture
def pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    left = slide.shapes.add_textbox(100, 100, 1000000, 300000)
    left.text = "금주 예정 업무"
    right = slide.shapes.add_textbox(1500000, 100, 1000000, 300000)
    right.text = "주간업무 실적사항"
    table = slide.shapes.add_table(2, 2, 100, 500000, 2500000, 600000).table
    table.cell(0, 0).text = "요일"
    table.cell(0, 1).text = "계획"
    table.cell(1, 0).text = "월"
    table.cell(1, 1).text = "정기회의"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
